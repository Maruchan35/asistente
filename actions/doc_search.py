# -*- coding: utf-8 -*-
"""
doc_search.py — Búsqueda DENTRO de los documentos del usuario (RAG local ligero).

Indexa el TEXTO de los documentos en Desktop / Documents / Downloads
(txt, md, docx, pdf, pptx) y permite preguntas tipo:
  "¿qué decía mi documento de métodos numéricos?"
  "¿en qué archivo hablo del presupuesto?"

Diseño (sin dependencias nuevas):
  • Índice invertido léxico con scoring TF-IDF simplificado
  • Cache en config/doc_index.json — solo re-lee archivos modificados
  • Extracción: python-docx (docx), python-pptx (pptx), pypdf si está / texto plano
  • Límites: max 500 archivos, 200KB de texto por archivo

La primera búsqueda construye el índice (~10-60s según cuántos documentos);
las siguientes son instantáneas con refresh incremental.
"""
from __future__ import annotations
import json
import logging
import math
import os
import re
import time
import unicodedata
from pathlib import Path

# pypdf es ruidoso con PDFs corruptos — silenciar (los errores ya se manejan)
logging.getLogger("pypdf").setLevel(logging.CRITICAL)
logging.getLogger("PyPDF2").setLevel(logging.CRITICAL)

_BASE       = Path(__file__).resolve().parent.parent
_INDEX_PATH = _BASE / "config" / "doc_index.json"
_HOME       = Path(os.path.expanduser("~"))

_SCAN_DIRS  = [_HOME / "Desktop", _HOME / "Documents", _HOME / "Downloads"]
_EXTENSIONS = {".txt", ".md", ".docx", ".pdf", ".pptx"}
_MAX_FILES        = 500
_MAX_TEXT_CHARS   = 200_000     # por archivo
_MAX_FILE_MB      = 40          # ignorar archivos gigantes
_SNIPPET_CHARS    = 240

# Carpetas a ignorar dentro de los directorios escaneados
_SKIP_DIR_NAMES = {".git", ".venv", "node_modules", "__pycache__", "backups",
                   "AppData", "$RECYCLE.BIN"}

_STOPWORDS = {
    "el", "la", "los", "las", "un", "una", "unos", "unas", "de", "del", "en",
    "y", "o", "a", "que", "es", "se", "por", "con", "para", "su", "sus", "al",
    "lo", "como", "más", "mas", "pero", "sí", "si", "no", "mi", "mis", "este",
    "esta", "esto", "ese", "esa", "eso", "the", "of", "and", "to", "in", "a",
    "is", "it", "on", "for", "qué", "cual", "cuál", "donde", "dónde",
}


# ══════════════════════════════════════════════════════════════════════════════
#  EXTRACCIÓN DE TEXTO
# ══════════════════════════════════════════════════════════════════════════════

def _extract_text(path: Path) -> str:
    """Extraer texto plano de un archivo según su tipo. '' si falla."""
    suffix = path.suffix.lower()
    try:
        if suffix in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="ignore")[:_MAX_TEXT_CHARS]

        if suffix == ".docx":
            from docx import Document
            doc = Document(str(path))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            # Incluir texto de tablas
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
            return "\n".join(parts)[:_MAX_TEXT_CHARS]

        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)[:_MAX_TEXT_CHARS]

        if suffix == ".pdf":
            # Intentar pypdf → PyPDF2 → nada
            try:
                from pypdf import PdfReader
            except ImportError:
                try:
                    from PyPDF2 import PdfReader
                except ImportError:
                    return ""
            reader = PdfReader(str(path))
            parts = []
            total = 0
            for page in reader.pages[:60]:
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                parts.append(t)
                total += len(t)
                if total > _MAX_TEXT_CHARS:
                    break
            return "\n".join(parts)[:_MAX_TEXT_CHARS]
    except Exception:
        return ""
    return ""


# ══════════════════════════════════════════════════════════════════════════════
#  TOKENIZACIÓN + ÍNDICE
# ══════════════════════════════════════════════════════════════════════════════

def _normalize(text: str) -> str:
    """Minúsculas + sin acentos (búsqueda tolerante)."""
    text = text.lower()
    text = unicodedata.normalize("NFD", text)
    return "".join(c for c in text if unicodedata.category(c) != "Mn")


def _tokenize(text: str) -> list[str]:
    words = re.findall(r"[a-z0-9ñ]{3,}", _normalize(text))
    return [w for w in words if w not in _STOPWORDS]


def _load_index() -> dict:
    try:
        if _INDEX_PATH.exists():
            return json.loads(_INDEX_PATH.read_text(encoding="utf-8"))
    except Exception:
        pass
    return {"files": {}, "built_at": 0}


def _save_index(index: dict) -> None:
    try:
        from core.safe_json import safe_write
        safe_write(_INDEX_PATH, index)
    except Exception:
        try:
            _INDEX_PATH.write_text(json.dumps(index, ensure_ascii=False),
                                   encoding="utf-8")
        except Exception:
            pass


def _scan_files() -> list[Path]:
    """Listar candidatos a indexar."""
    out: list[Path] = []
    for root_dir in _SCAN_DIRS:
        if not root_dir.is_dir():
            continue
        for dirpath, dirnames, filenames in os.walk(root_dir):
            # Podar carpetas ignoradas
            dirnames[:] = [d for d in dirnames
                           if d not in _SKIP_DIR_NAMES and not d.startswith(".")]
            for fn in filenames:
                p = Path(dirpath) / fn
                if p.suffix.lower() not in _EXTENSIONS:
                    continue
                try:
                    if p.stat().st_size > _MAX_FILE_MB * 1024 * 1024:
                        continue
                except OSError:
                    continue
                out.append(p)
                if len(out) >= _MAX_FILES:
                    return out
    return out


def _refresh_index(force: bool = False, progress_cb=None) -> dict:
    """Actualizar el índice: solo re-procesa archivos nuevos o modificados."""
    index = {"files": {}, "built_at": 0} if force else _load_index()
    files_meta: dict = index.get("files", {})
    current = _scan_files()
    current_keys = set()
    changed = 0

    for i, p in enumerate(current):
        key = str(p)
        current_keys.add(key)
        try:
            mtime = p.stat().st_mtime
        except OSError:
            continue
        meta = files_meta.get(key)
        if meta and abs(meta.get("mtime", 0) - mtime) < 1:
            continue   # sin cambios
        if progress_cb and changed % 10 == 0:
            progress_cb(f"Indexando {i+1}/{len(current)}: {p.name[:40]}")
        text = _extract_text(p)
        if not text.strip():
            files_meta.pop(key, None)
            continue
        tokens = _tokenize(text)
        if not tokens:
            continue
        # Term frequencies (top 400 términos por doc para acotar tamaño)
        tf: dict[str, int] = {}
        for t in tokens:
            tf[t] = tf.get(t, 0) + 1
        top_tf = dict(sorted(tf.items(), key=lambda kv: -kv[1])[:400])
        files_meta[key] = {
            "mtime":   mtime,
            "name":    p.name,
            "tf":      top_tf,
            "n_tokens": len(tokens),
            "preview": text[:_SNIPPET_CHARS * 4],
        }
        changed += 1

    # Eliminar archivos que ya no existen
    for key in list(files_meta.keys()):
        if key not in current_keys:
            del files_meta[key]

    index["files"] = files_meta
    index["built_at"] = time.time()
    if changed or force:
        _save_index(index)
    return index


# ══════════════════════════════════════════════════════════════════════════════
#  BÚSQUEDA
# ══════════════════════════════════════════════════════════════════════════════

def _search(query: str, index: dict, max_results: int = 5) -> list[dict]:
    """TF-IDF score sobre el índice."""
    q_tokens = _tokenize(query)
    if not q_tokens:
        return []
    files = index.get("files", {})
    n_docs = max(1, len(files))

    # Document frequency por término de la query
    df: dict[str, int] = {}
    for t in q_tokens:
        df[t] = sum(1 for meta in files.values() if t in meta.get("tf", {}))

    scored = []
    for path, meta in files.items():
        tf = meta.get("tf", {})
        n_tokens = max(1, meta.get("n_tokens", 1))
        score = 0.0
        matched_terms = []
        for t in q_tokens:
            if t in tf:
                idf = math.log(1 + n_docs / max(1, df.get(t, 1)))
                score += (tf[t] / n_tokens) * idf
                matched_terms.append(t)
        # Bonus por nombre de archivo que matchea
        name_norm = _normalize(meta.get("name", ""))
        for t in q_tokens:
            if t in name_norm:
                score += 0.05
                if t not in matched_terms:
                    matched_terms.append(t)
        if score > 0:
            scored.append((score, path, meta, matched_terms))

    scored.sort(key=lambda x: -x[0])
    out = []
    for score, path, meta, matched in scored[:max_results]:
        # Snippet: buscar primer término en el preview
        preview = meta.get("preview", "")
        snippet = preview[:_SNIPPET_CHARS]
        pv_norm = _normalize(preview)
        for t in matched:
            pos = pv_norm.find(t)
            if pos > 40:
                snippet = "…" + preview[max(0, pos-60):pos + _SNIPPET_CHARS - 60]
                break
        out.append({
            "path":    path,
            "name":    meta.get("name", ""),
            "score":   round(score, 4),
            "matched": matched,
            "snippet": snippet.replace("\n", " ").strip(),
        })
    return out


# ══════════════════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

def doc_search(parameters: dict, player=None) -> str:
    action      = (parameters.get("action") or "search").lower()
    query       = (parameters.get("query") or "").strip()
    max_results = int(parameters.get("max_results", 5))

    def _log(msg):
        if player:
            try: player.write_log(msg)
            except Exception: pass

    if action == "reindex":
        _log("📚 Reindexando documentos (esto puede tardar)...")
        index = _refresh_index(force=True, progress_cb=_log)
        n = len(index.get("files", {}))
        return f"Índice reconstruido: {n} documentos indexados."

    if action == "stats":
        index = _load_index()
        n = len(index.get("files", {}))
        when = index.get("built_at", 0)
        age = f"{int((time.time()-when)/60)} min" if when else "nunca"
        return f"Índice: {n} documentos, última actualización hace {age}."

    # search (default)
    if not query:
        return "Error: falta el parámetro 'query'."

    index = _load_index()
    if not index.get("files"):
        _log("📚 Primera búsqueda — construyendo índice de documentos...")
        index = _refresh_index(force=False, progress_cb=_log)
    else:
        # Refresh incremental rápido (solo cambios)
        index = _refresh_index(force=False)

    results = _search(query, index, max_results=max_results)
    if not results:
        n = len(index.get("files", {}))
        return (f"No encontré documentos sobre '{query}' "
                f"(busqué en {n} documentos de Desktop/Documents/Downloads).")

    lines = [f"Documentos sobre '{query}':"]
    for i, r in enumerate(results, 1):
        lines.append(f"\n{i}. {r['name']}")
        lines.append(f"   Ruta: {r['path']}")
        if r["snippet"]:
            lines.append(f"   Extracto: \"{r['snippet']}\"")
    lines.append(
        "\n(Para leer el contenido completo de uno, usa file_processor "
        "con la ruta indicada.)"
    )
    return "\n".join(lines)
